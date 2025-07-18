#!/usr/bin/env python3
"""
PowerPoint generation service for creating presentations.

This module handles the core PowerPoint generation functionality using python-pptx,
including slide creation, theme application, chart embedding, and file export.
"""

import asyncio
import json
import os
import time
from datetime import datetime, timedelta
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

import aiofiles
import httpx
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from structlog import get_logger

from app.core.config import settings
from app.core.database import execute_query, execute_transaction
from app.core.redis_client import redis_manager, CacheManager
from app.models.powerpoint_models import (
    JobStatus,
    OutputFormat,
    Theme,
    SlideType,
    LayoutType,
    AnimationType,
    TransitionType,
    ChartType,
    ColorScheme,
    PresentationConfig,
    Slide,
    Chart,
    TextContent,
    ImageContent,
    TableContent
)


logger = get_logger(__name__)


class PowerPointGenerator:
    """PowerPoint generation service."""
    
    def __init__(self):
        self.cache = CacheManager()
        self.active_jobs: Dict[int, Dict[str, Any]] = {}
        self._lock = asyncio.Lock()
    
    async def generate_presentation(
        self,
        job_id: int,
        user_id: int,
        presentation_config: PresentationConfig,
        output_format: OutputFormat,
        data_source: Dict[str, Any]
    ) -> Tuple[bool, Optional[str], Optional[str]]:
        """Generate PowerPoint presentation.
        
        Returns:
            Tuple of (success, file_path, error_message)
        """
        start_time = time.time()
        
        try:
            # Update job status to processing
            await self._update_job_status(job_id, JobStatus.PROCESSING, started_at=datetime.utcnow())
            
            # Track active job
            async with self._lock:
                self.active_jobs[job_id] = {
                    "status": JobStatus.PROCESSING,
                    "start_time": start_time,
                    "current_slide": 0,
                    "total_slides": len(presentation_config.slides)
                }
            
            logger.info("Starting PowerPoint generation", job_id=job_id, user_id=user_id)
            
            # Fetch data if needed
            data = await self._fetch_data(data_source)
            
            # Create presentation
            prs = await self._create_presentation(presentation_config, data, job_id)
            
            # Generate file path
            file_name = f"presentation_{job_id}_{int(time.time())}.{output_format.value}"
            file_path = os.path.join(settings.PPT_OUTPUT_DIR, file_name)
            
            # Save presentation
            if output_format == OutputFormat.PPTX:
                prs.save(file_path)
            else:
                # Convert to other formats
                await self._convert_presentation(prs, file_path, output_format)
            
            # Get file stats
            file_size = os.path.getsize(file_path)
            slide_count = len(prs.slides)
            chart_count = await self._count_charts(presentation_config)
            animation_count = await self._count_animations(presentation_config)
            
            # Calculate generation time
            generation_time_ms = int((time.time() - start_time) * 1000)
            
            # Update job with success
            await self._update_job_completion(
                job_id=job_id,
                status=JobStatus.COMPLETED,
                file_path=file_path,
                file_size=file_size,
                slide_count=slide_count,
                chart_count=chart_count,
                animation_count=animation_count,
                generation_time_ms=generation_time_ms
            )
            
            # Clean up active job tracking
            async with self._lock:
                self.active_jobs.pop(job_id, None)
            
            logger.info(
                "PowerPoint generation completed",
                job_id=job_id,
                file_path=file_path,
                file_size=file_size,
                generation_time_ms=generation_time_ms
            )
            
            return True, file_path, None
        
        except Exception as e:
            error_message = str(e)
            logger.error("PowerPoint generation failed", job_id=job_id, error=error_message)
            
            # Update job with failure
            await self._update_job_status(
                job_id,
                JobStatus.FAILED,
                error_message=error_message,
                completed_at=datetime.utcnow()
            )
            
            # Clean up active job tracking
            async with self._lock:
                self.active_jobs.pop(job_id, None)
            
            return False, None, error_message
    
    async def _create_presentation(
        self,
        config: PresentationConfig,
        data: Dict[str, Any],
        job_id: int
    ) -> Presentation:
        """Create PowerPoint presentation from configuration."""
        # Create new presentation
        prs = Presentation()
        
        # Set slide size
        prs.slide_width = Inches(config.slide_size.width)
        prs.slide_height = Inches(config.slide_size.height)
        
        # Apply theme
        await self._apply_theme(prs, config.theme, config.color_scheme)
        
        # Create slides
        for i, slide_config in enumerate(config.slides):
            await self._update_job_progress(job_id, i, len(config.slides))
            
            slide = await self._create_slide(prs, slide_config, data, config)
            
            # Add slide notes if specified
            if slide_config.notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_config.notes
        
        # Set presentation metadata
        await self._set_presentation_metadata(prs, config.metadata)
        
        return prs
    
    async def _create_slide(
        self,
        prs: Presentation,
        slide_config: Slide,
        data: Dict[str, Any],
        presentation_config: PresentationConfig
    ) -> Any:
        """Create a single slide."""
        # Get slide layout
        layout_index = self._get_layout_index(slide_config.layout)
        slide_layout = prs.slide_layouts[layout_index]
        
        # Add slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background
        await self._set_slide_background(slide, slide_config)
        
        # Add title if specified
        if slide_config.title and hasattr(slide, 'shapes') and slide.shapes.title:
            slide.shapes.title.text = slide_config.title
            await self._apply_text_style(slide.shapes.title.text_frame, presentation_config.default_font)
        
        # Add content elements
        await self._add_slide_content(slide, slide_config, data, presentation_config)
        
        return slide
    
    async def _add_slide_content(
        self,
        slide: Any,
        slide_config: Slide,
        data: Dict[str, Any],
        presentation_config: PresentationConfig
    ) -> None:
        """Add content elements to slide."""
        content = slide_config.content
        
        # Add text elements
        for text_element in content.texts:
            await self._add_text_element(slide, text_element)
        
        # Add image elements
        for image_element in content.images:
            await self._add_image_element(slide, image_element)
        
        # Add chart elements
        for chart_element in content.charts:
            await self._add_chart_element(slide, chart_element, data)
        
        # Add table elements
        for table_element in content.tables:
            await self._add_table_element(slide, table_element, data)
    
    async def _add_text_element(self, slide: Any, text_element: TextContent) -> None:
        """Add text element to slide."""
        # Convert position to inches
        left = Inches(text_element.position.x)
        top = Inches(text_element.position.y)
        width = Inches(text_element.position.width)
        height = Inches(text_element.position.height)
        
        # Add text box
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text_element.text
        
        # Apply text styling
        await self._apply_text_style(text_frame, text_element.style.font)
        
        # Set alignment
        for paragraph in text_frame.paragraphs:
            alignment = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY
            }.get(text_element.style.alignment, PP_ALIGN.LEFT)
            paragraph.alignment = alignment
    
    async def _add_image_element(self, slide: Any, image_element: ImageContent) -> None:
        """Add image element to slide."""
        try:
            # Download or load image
            image_data = await self._load_image(image_element.image_url)
            
            # Convert position to inches
            left = Inches(image_element.position.x)
            top = Inches(image_element.position.y)
            width = Inches(image_element.position.width)
            height = Inches(image_element.position.height)
            
            # Add image to slide
            pic = slide.shapes.add_picture(BytesIO(image_data), left, top, width, height)
            
            # Set alt text if provided
            if image_element.alt_text:
                pic.element.set("title", image_element.alt_text)
        
        except Exception as e:
            logger.error("Failed to add image element", error=str(e))
    
    async def _add_chart_element(self, slide: Any, chart_element: Chart, data: Dict[str, Any]) -> None:
        """Add chart element to slide."""
        try:
            # Prepare chart data
            chart_data = CategoryChartData()
            chart_data.categories = chart_element.data.labels
            
            for dataset in chart_element.data.datasets:
                chart_data.add_series(dataset["label"], dataset["data"])
            
            # Convert position to inches
            left = Inches(chart_element.position.x)
            top = Inches(chart_element.position.y)
            width = Inches(chart_element.position.width)
            height = Inches(chart_element.position.height)
            
            # Get chart type
            chart_type = self._get_chart_type(chart_element.config.chart_type)
            
            # Add chart to slide
            chart = slide.shapes.add_chart(
                chart_type, left, top, width, height, chart_data
            ).chart
            
            # Configure chart
            await self._configure_chart(chart, chart_element.config)
        
        except Exception as e:
            logger.error("Failed to add chart element", error=str(e))
    
    async def _add_table_element(self, slide: Any, table_element: TableContent, data: Dict[str, Any]) -> None:
        """Add table element to slide."""
        try:
            # Calculate table dimensions
            rows = len(table_element.rows) + 1  # +1 for header
            cols = len(table_element.headers)
            
            # Convert position to inches
            left = Inches(table_element.position.x)
            top = Inches(table_element.position.y)
            width = Inches(table_element.position.width)
            height = Inches(table_element.position.height)
            
            # Add table to slide
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
            
            # Set headers
            for i, header in enumerate(table_element.headers):
                cell = table.cell(0, i)
                cell.text = header
                await self._apply_table_cell_style(cell, table_element.header_style)
            
            # Set rows
            for row_idx, row_data in enumerate(table_element.rows):
                for col_idx, cell_data in enumerate(row_data.cells):
                    if col_idx < cols:
                        cell = table.cell(row_idx + 1, col_idx)
                        cell.text = str(cell_data)
                        style = row_data.style or table_element.row_style
                        await self._apply_table_cell_style(cell, style)
        
        except Exception as e:
            logger.error("Failed to add table element", error=str(e))
    
    async def _apply_theme(self, prs: Presentation, theme: Theme, color_scheme: ColorScheme) -> None:
        """Apply theme to presentation."""
        # Theme application logic would go here
        # This is a simplified implementation
        logger.info("Applying theme", theme=theme.value, color_scheme=color_scheme.value)
    
    async def _apply_text_style(self, text_frame: Any, font_style: Any) -> None:
        """Apply text styling to text frame."""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.name = font_style.family
                font.size = Pt(font_style.size)
                font.bold = font_style.bold
                font.italic = font_style.italic
                font.underline = font_style.underline
                
                # Set color
                if font_style.color.startswith("#"):
                    color_hex = font_style.color[1:]
                    r = int(color_hex[0:2], 16)
                    g = int(color_hex[2:4], 16)
                    b = int(color_hex[4:6], 16)
                    font.color.rgb = RGBColor(r, g, b)
    
    async def _apply_table_cell_style(self, cell: Any, style: Any) -> None:
        """Apply styling to table cell."""
        if style:
            await self._apply_text_style(cell.text_frame, style.font)
    
    async def _set_slide_background(self, slide: Any, slide_config: Slide) -> None:
        """Set slide background."""
        # Background setting logic would go here
        # This is a simplified implementation
        pass
    
    async def _configure_chart(self, chart: Any, config: Any) -> None:
        """Configure chart properties."""
        # Chart configuration logic would go here
        # This is a simplified implementation
        if config.title:
            chart.chart_title.text_frame.text = config.title
        
        chart.has_legend = config.show_legend
    
    def _get_layout_index(self, layout_type: LayoutType) -> int:
        """Get slide layout index for layout type."""
        layout_mapping = {
            LayoutType.TITLE_SLIDE: 0,
            LayoutType.TITLE_AND_CONTENT: 1,
            LayoutType.TWO_CONTENT: 3,
            LayoutType.COMPARISON: 4,
            LayoutType.TITLE_ONLY: 5,
            LayoutType.BLANK: 6,
            LayoutType.CONTENT_WITH_CAPTION: 7,
            LayoutType.PICTURE_WITH_CAPTION: 8
        }
        return layout_mapping.get(layout_type, 1)  # Default to title and content
    
    def _get_chart_type(self, chart_type: ChartType) -> Any:
        """Get python-pptx chart type from our enum."""
        chart_mapping = {
            ChartType.BAR: XL_CHART_TYPE.BAR_CLUSTERED,
            ChartType.COLUMN: XL_CHART_TYPE.COLUMN_CLUSTERED,
            ChartType.LINE: XL_CHART_TYPE.LINE,
            ChartType.PIE: XL_CHART_TYPE.PIE,
            ChartType.AREA: XL_CHART_TYPE.AREA,
            ChartType.SCATTER: XL_CHART_TYPE.XY_SCATTER,
            ChartType.DOUGHNUT: XL_CHART_TYPE.DOUGHNUT,
            ChartType.RADAR: XL_CHART_TYPE.RADAR
        }
        return chart_mapping.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    async def _load_image(self, image_url: str) -> bytes:
        """Load image from URL or file path."""
        if image_url.startswith("http"):
            # Download from URL
            async with httpx.AsyncClient() as client:
                response = await client.get(image_url, timeout=30)
                response.raise_for_status()
                return response.content
        else:
            # Load from file path
            async with aiofiles.open(image_url, "rb") as f:
                return await f.read()
    
    async def _fetch_data(self, data_source: Dict[str, Any]) -> Dict[str, Any]:
        """Fetch data from data source."""
        # Data fetching logic would go here
        # This is a simplified implementation
        return data_source.get("static_source", {}).get("data", {})
    
    async def _convert_presentation(self, prs: Presentation, file_path: str, output_format: OutputFormat) -> None:
        """Convert presentation to different formats."""
        # Conversion logic would go here
        # This is a simplified implementation that just saves as PPTX
        prs.save(file_path)
    
    async def _set_presentation_metadata(self, prs: Presentation, metadata: Any) -> None:
        """Set presentation metadata."""
        core_props = prs.core_properties
        core_props.title = metadata.title
        core_props.author = metadata.author
        core_props.subject = metadata.subject
        core_props.keywords = ", ".join(metadata.keywords)
        core_props.category = metadata.category
        core_props.comments = metadata.description
    
    async def _count_charts(self, config: PresentationConfig) -> int:
        """Count total charts in presentation."""
        return sum(len(slide.content.charts) for slide in config.slides)
    
    async def _count_animations(self, config: PresentationConfig) -> int:
        """Count total animations in presentation."""
        return sum(1 for slide in config.slides if slide.animation != AnimationType.NONE)
    
    async def _update_job_status(
        self,
        job_id: int,
        status: JobStatus,
        error_message: Optional[str] = None,
        started_at: Optional[datetime] = None,
        completed_at: Optional[datetime] = None
    ) -> None:
        """Update job status in database."""
        update_fields = ["status = $2"]
        params = [job_id, status.value]
        param_count = 2
        
        if error_message is not None:
            param_count += 1
            update_fields.append(f"error_message = ${param_count}")
            params.append(error_message)
        
        if started_at is not None:
            param_count += 1
            update_fields.append(f"started_at = ${param_count}")
            params.append(started_at)
        
        if completed_at is not None:
            param_count += 1
            update_fields.append(f"completed_at = ${param_count}")
            params.append(completed_at)
        
        query = f"""
            UPDATE ppt_export_jobs 
            SET {', '.join(update_fields)}
            WHERE id = $1
        """
        
        await execute_query(query, *params, fetch="none")
    
    async def _update_job_completion(
        self,
        job_id: int,
        status: JobStatus,
        file_path: str,
        file_size: int,
        slide_count: int,
        chart_count: int,
        animation_count: int,
        generation_time_ms: int
    ) -> None:
        """Update job with completion details."""
        query = """
            UPDATE ppt_export_jobs 
            SET status = $2, file_path = $3, file_size = $4, slide_count = $5,
                chart_count = $6, animation_count = $7, generation_time_ms = $8,
                completed_at = $9
            WHERE id = $1
        """
        
        await execute_query(
            query,
            job_id,
            status.value,
            file_path,
            file_size,
            slide_count,
            chart_count,
            animation_count,
            generation_time_ms,
            datetime.utcnow(),
            fetch="none"
        )
    
    async def _update_job_progress(self, job_id: int, current_slide: int, total_slides: int) -> None:
        """Update job progress."""
        async with self._lock:
            if job_id in self.active_jobs:
                self.active_jobs[job_id]["current_slide"] = current_slide
                self.active_jobs[job_id]["total_slides"] = total_slides
    
    async def get_job_status(self, job_id: int) -> Optional[Dict[str, Any]]:
        """Get job status and progress."""
        async with self._lock:
            if job_id in self.active_jobs:
                job_info = self.active_jobs[job_id]
                runtime = time.time() - job_info["start_time"]
                progress = (job_info["current_slide"] / job_info["total_slides"]) * 100 if job_info["total_slides"] > 0 else 0
                
                return {
                    "job_id": job_id,
                    "status": job_info["status"].value,
                    "progress_percentage": round(progress, 2),
                    "current_slide": job_info["current_slide"],
                    "total_slides": job_info["total_slides"],
                    "runtime_seconds": round(runtime, 2)
                }
        
        return None
    
    async def cancel_job(self, job_id: int) -> bool:
        """Cancel a running job."""
        async with self._lock:
            if job_id in self.active_jobs:
                self.active_jobs[job_id]["status"] = JobStatus.CANCELLED
                
                # Update database
                await self._update_job_status(
                    job_id,
                    JobStatus.CANCELLED,
                    completed_at=datetime.utcnow()
                )
                
                return True
        
        return False


# Global generator instance
powerpoint_generator = PowerPointGenerator()


# Export commonly used functions
__all__ = ["powerpoint_generator", "PowerPointGenerator"]
